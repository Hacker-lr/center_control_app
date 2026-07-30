#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
中控系统 App 功能介绍 PPT 生成脚本
生成 .pptx 格式的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏比例
prs.slide_height = Inches(7.5)

# 定义颜色（与 App 一致）
BG_COLOR = RGBColor(0x0D, 0x11, 0x17)       # 深色背景
TEXT_COLOR = RGBColor(0xD4, 0xC5, 0xA9)     # 主文字颜色
ACCENT_COLOR = RGBColor(0x1F, 0x40, 0x68)   # 主色调
HIGHLIGHT_COLOR = RGBColor(0x3E, 0x6B, 0x48) # 高亮色
SECONDARY_COLOR = RGBColor(0x8B, 0x94, 0x9E) # 次要文字颜色

def set_slide_background(slide):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, title_text, font_size=44, top=0.5):
    """添加标题"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.333), Inches(1))
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].text = title_text
    title_frame.paragraphs[0].font.size = Pt(font_size)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_subtitle(slide, text, top=1.3):
    """添加副标题"""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.333), Inches(0.6))
    frame = box.text_frame
    frame.paragraphs[0].text = text
    frame.paragraphs[0].font.size = Pt(18)
    frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_bullet_list(slide, items, left=1, top=2, width=11, font_size=20, spacing=0.4):
    """添加项目列表"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(font_size * spacing)
        p.level = 0

def add_card(slide, title, description, left, top, width=5.5, height=1.5):
    """添加卡片"""
    # 卡片背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
    shape.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.2), Inches(width - 0.6), Inches(0.4))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_COLOR
    
    # 描述
    desc_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.6), Inches(width - 0.6), Inches(height - 0.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = description
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR

# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_background(slide1)

# 主标题
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
tf = title_box.text_frame
tf.paragraphs[0].text = "中控系统 App"
tf.paragraphs[0].font.size = Pt(56)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT_COLOR
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 副标题
sub_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.6))
tf = sub_box.text_frame
tf.paragraphs[0].text = "基于 Flutter 的跨平台设备集中控制系统"
tf.paragraphs[0].font.size = Pt(24)
tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 版本信息
ver_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.5))
tf = ver_box.text_frame
tf.paragraphs[0].text = "v1.0.0 | 2026"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# 第2页：系统概述
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2)
add_title(slide2, "系统概述")

features = [
    ("跨平台", "基于 Flutter 开发，一套代码同时支持 Android、iOS、Web 及桌面端"),
    ("双协议通信", "支持 TCP 和 UDP 两种网络协议，每种设备可独立选择"),
    ("品牌配置", "内置多种设备品牌模板，选择品牌后自动填充协议、端口和指令"),
    ("持久化存储", "所有配置参数自动保存到本地，重启后保持"),
]

positions = [(0.8, 2), (6.8, 2), (0.8, 4), (6.8, 4)]
for (title, desc), (left, top) in zip(features, positions):
    add_card(slide2, title, desc, left, top, 5.5, 1.6)

# ============================================================
# 第3页：四大设备控制
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3)
add_title(slide3, "四大设备控制")

devices = [
    "⚡ 时序电源控制：远程开关、TCP/UDP 双协议、状态反馈",
    "📺 大屏拼接器控制：分屏切换、可视化预览、输入绑定",
    "🎬 视频矩阵控制：输入/输出绑定、长按改名、分页显示",
    "📹 摄像头控制：Sony VISCA 协议、云台控制、预置位管理",
]

add_bullet_list(slide3, devices, left=1.5, top=2.2, font_size=24, spacing=0.6)

# ============================================================
# 第4页：核心交互特性
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4)
add_title(slide4, "核心交互特性")

features = [
    ("长按改名", "视频矩阵、摄像头按钮均支持长按 2 秒触发重命名"),
    ("实时状态同步", "多页面共享矩阵状态，通过 ChangeNotifier 实时同步"),
    ("自动重连", "设备断线后自动重连，带心跳检测机制"),
    ("响应式布局", "自适应手机、平板、桌面设备，支持横竖屏切换"),
]

positions = [(0.8, 2), (6.8, 2), (0.8, 4), (6.8, 4)]
for (title, desc), (left, top) in zip(features, positions):
    add_card(slide4, title, desc, left, top, 5.5, 1.6)

# ============================================================
# 第5页：调试配置系统
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5)
add_title(slide5, "调试配置系统")
add_subtitle(slide5, '长按顶部标题"欢迎使用中控"进入配置页面', top=1.2)

configs = [
    "1. 设备参数配置：修改所有受控设备的 IP 地址、端口号、通信协议",
    "2. 品牌选择：选择设备品牌后自动填充协议、端口和控制指令",
    "3. 通道数量设置：自定义矩阵输入/输出通道数、摄像机个数、预置位数量",
    "4. 页面显示开关：控制各功能页面的显示/隐藏，按需启用控制模块",
]

add_bullet_list(slide5, configs, left=1.5, top=2.2, font_size=20, spacing=0.5)

# ============================================================
# 第6页：技术架构
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide6)
add_title(slide6, "技术架构")

# 技术栈
tech_items = [
    ("Flutter", "跨平台 UI 框架\nDart 语言开发"),
    ("TCP/UDP", "双协议网络通信\n支持 ASCII/Hex"),
    ("SharedPreferences", "本地持久化存储\n配置自动保存"),
]

for i, (name, desc) in enumerate(tech_items):
    left = 1.5 + i * 4
    
    # 技术项框
    shape = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(2.2), Inches(3.2), Inches(1.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
    shape.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    
    # 名称
    name_box = slide6.shapes.add_textbox(Inches(left), Inches(2.4), Inches(3.2), Inches(0.5))
    tf = name_box.text_frame
    tf.paragraphs[0].text = name
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_COLOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 描述
    desc_box = slide6.shapes.add_textbox(Inches(left), Inches(3), Inches(3.2), Inches(0.8))
    tf = desc_box.text_frame
    tf.paragraphs[0].text = desc
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 架构流程
flow_items = ["Flutter App", "DeviceConfig", "BaseConnection", "硬件设备"]
flow_descs = ["UI 交互层", "配置中心", "网络连接层", "电源/矩阵/大屏"]

for i, (item, desc) in enumerate(zip(flow_items, flow_descs)):
    left = 1 + i * 3.2
    
    shape = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(4.8), Inches(2.4), Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1F, 0x40, 0x68)
    shape.line.color.rgb = ACCENT_COLOR
    
    box = slide6.shapes.add_textbox(Inches(left), Inches(4.9), Inches(2.4), Inches(0.8))
    tf = box.text_frame
    tf.paragraphs[0].text = item
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_COLOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 箭头
    if i < len(flow_items) - 1:
        arrow_box = slide6.shapes.add_textbox(Inches(left + 2.5), Inches(5.1), Inches(0.6), Inches(0.4))
        tf = arrow_box.text_frame
        tf.paragraphs[0].text = "→"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# 第7页：功能清单
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide7)
add_title(slide7, "功能清单")

features = [
    "时序电源远程开关控制（TCP/UDP 双协议）",
    "大屏拼接器分屏模式切换（全屏/二分/三分/四分/五分）",
    "视频矩阵输入/输出通道绑定与分页显示",
    "摄像头云台控制与预置位管理（Sony VISCA 协议）",
    "设备品牌配置（协议/端口/指令自动填充）",
    "页面显示开关控制（按需启用功能模块）",
    "通道/预置位长按改名（文字自适应缩放）",
    "调试配置页面（IP/端口/通道数/预置位数）",
    "自动重连与心跳检测（断线自动恢复）",
    "响应式布局（自适应手机/平板/桌面）",
]

add_bullet_list(slide7, features, left=1.2, top=1.8, font_size=18, spacing=0.35)

# ============================================================
# 第8页：结束页
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide8)

# 感谢观看
title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
tf = title_box.text_frame
tf.paragraphs[0].text = "感谢观看"
tf.paragraphs[0].font.size = Pt(48)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT_COLOR
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 副标题
sub_box = slide8.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.6))
tf = sub_box.text_frame
tf.paragraphs[0].text = "中控系统 App v1.0.0"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 技术栈
tech_box = slide8.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.8))
tf = tech_box.text_frame
tf.paragraphs[0].text = "技术栈：Flutter + Dart + TCP/UDP + SharedPreferences"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 保存
output_path = "c:/Users/YSL/Desktop/参考/center_control_app/presentation/中控系统App功能介绍.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")